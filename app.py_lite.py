import sys
import os
import traceback
from PyQt5.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QHBoxLayout, 
    QPushButton, QLineEdit, QLabel, QListWidget, 
    QFileDialog, QMessageBox, QSizePolicy, QProgressDialog
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, QCoreApplication
from pptx import Presentation

# win32com 관련 모듈을 Windows 환경에서만 import하도록 처리
try:
    import pythoncom
    import win32com.client
    WINDOWS_ENV = True
except ImportError:
    WINDOWS_ENV = False

# PPTX 병합 작업을 별도의 스레드에서 처리하기 위한 워커 클래스
class MergerWorker(QThread):
    # 작업 진행 상황을 GUI에 알리기 위한 시그널
    merge_finished = pyqtSignal(bool, str) # (성공 여부, 메시지)
    progress_update = pyqtSignal(int, int) # (현재 슬라이드 번호, 전체 슬라이드 개수)

    def __init__(self, file_paths, output_path, parent=None):
        super().__init__(parent)
        self.file_paths = file_paths
        self.output_path = output_path
        self.temp_files = [] # 변환된 임시 파일 목록

    def __del__(self):
        # 스레드 종료 시 임시 파일 정리
        for temp_file in self.temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception:
                pass

    # .ppt 파일을 .pptx로 변환하는 함수 (Windows + MS PowerPoint 필요)
    def convert_ppt_to_pptx(self, ppt_path):
        if not WINDOWS_ENV:
            raise EnvironmentError("`.ppt` 파일 변환은 Windows 환경과 `pywin32` 라이브러리, 그리고 MS PowerPoint가 필요합니다.")
        
        pythoncom.CoInitialize() 
        
        temp_dir = os.path.dirname(ppt_path)
        temp_pptx_path = os.path.join(temp_dir, f"~temp_converted_{os.path.basename(ppt_path)[:-4]}.pptx")
        
        powerpoint = None
        try:
            powerpoint = win32com.client.Dispatch("Powerpoint.Application")
            powerpoint.Visible = 0

            presentation = powerpoint.Presentations.Open(
                ppt_path, 
                ReadOnly=True, 
                WithWindow=False
            )
            presentation.SaveAs(temp_pptx_path, 24) # 24는 ppSaveAsPresentation (pptx)
            presentation.Close()
            
            self.temp_files.append(temp_pptx_path)
            return temp_pptx_path
        except Exception as e:
            if powerpoint:
                try: powerpoint.Quit()
                except: pass
            raise Exception(f"PPT 파일 변환 실패 (MS PowerPoint 설치 및 권한 확인 필요): {e}")
        finally:
            if powerpoint:
                try: powerpoint.Quit()
                except: pass
            pythoncom.CoUninitialize()

    def run(self):
        if not self.file_paths:
            self.merge_finished.emit(False, "오류: 병합할 PPT 파일이 없습니다.")
            return
            
        process_paths = []
        try:
            # 1. 파일 목록 순회 및 PPTX로 변환
            for path in self.file_paths:
                if path.lower().endswith('.ppt'):
                    self.progress_update.emit(1, 1) # 변환 단계 표시
                    converted_path = self.convert_ppt_to_pptx(path)
                    process_paths.append(converted_path)
                elif path.lower().endswith('.pptx'):
                    process_paths.append(path)
                
            if not process_paths:
                self.merge_finished.emit(False, "오류: 병합할 파일이 유효하지 않습니다.")
                return

            # 2. 병합 로직 시작 및 슬라이드 카운트 최적화
            master_pptx = Presentation(process_paths[0])
            
            # --- 수정된 부분: len() 사용 ---
            total_slides_processed = len(master_pptx.slides) 
            total_slides_count = total_slides_processed
            
            # 나머지 파일들의 슬라이드 수만 합산하여 전체 카운트 계산
            for path in process_paths[1:]:
                # 파일을 로드하여 슬라이드 개수만 얻어옴
                total_slides_count += len(Presentation(path).slides) 

            self.progress_update.emit(total_slides_processed, total_slides_count)
            
            # 3. 나머지 파일들을 순회하며 슬라이드 복사
            for path in process_paths[1:]:
                source_pptx = Presentation(path)
                slide_layout_map = {layout.name: layout for layout in master_pptx.slide_layouts}

                for slide in source_pptx.slides:
                    source_layout_name = slide.slide_layout.name
                    target_layout = slide_layout_map.get(source_layout_name, master_pptx.slide_layouts[6])
                    
                    new_slide = master_pptx.slides.add_slide(target_layout)
                    
                    # 콘텐츠 복사 (텍스트만)
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text_frame = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height).text_frame
                            text_frame.text = shape.text
                        elif shape.shape_type == 13: 
                            pass # 이미지/차트 생략

                    total_slides_processed += 1
                    self.progress_update.emit(total_slides_processed, total_slides_count)
                    QCoreApplication.processEvents()

            # 4. 결과 파일 저장
            master_pptx.save(self.output_path)
            
            self.merge_finished.emit(True, f"✅ PPTX 병합 완료! (임시 파일 포함)\n\n저장 위치: {self.output_path}")

        except Exception as e:
            error_message = f"PPTX 병합 중 오류가 발생했습니다.\n\n오류: {e}\n\n상세:\n{traceback.format_exc()}"
            self.merge_finished.emit(False, error_message)
        finally:
            # 임시 파일 정리
            for temp_file in self.temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except Exception:
                    pass

class PptxMergerApp(QWidget):
    def __init__(self):
        super().__init__()
        title_suffix = " (ppt 자동 변환 기능 포함)" if WINDOWS_ENV else " (pptx 전용)"
        self.setWindowTitle('PPTX 순서 병합 프로그램' + title_suffix) 
        self.setGeometry(100, 100, 650, 480)
        self.save_directory = os.path.expanduser("~") 
        self.setup_ui()
        
        self.setAcceptDrops(True)
        self.worker = None

    def setup_ui(self):
        main_layout = QVBoxLayout()
        
        # --- 1. 파일 추가 및 드래그 앤 드롭 영역 ---
        file_select_layout = QHBoxLayout()
        
        self.select_files_button = QPushButton('+ 파일 추가 (탐색기)')
        self.select_files_button.clicked.connect(self.open_file_dialog)
        file_select_layout.addWidget(self.select_files_button)

        drag_info = "여기에 .PPTX 또는 .PPT 파일을 드래그 & 드롭 가능"
        if WINDOWS_ENV:
            drag_info += "\n(.PPT 파일은 자동으로 .PPTX로 변환됩니다. *MS PowerPoint 설치 필수)"
        
        drag_label = QLabel(drag_info)
        drag_label.setAlignment(Qt.AlignCenter)
        drag_label.setStyleSheet("border: 2px dashed #ccc; padding: 10px; color: #555; background-color: #f9f9f9; border-radius: 8px;")
        file_select_layout.addWidget(drag_label)
        
        main_layout.addLayout(file_select_layout)
        
        # --- 2. 병합 목록 영역 ---
        merge_list_group = QHBoxLayout()
        
        self.list_widget = QListWidget()
        self.list_widget.setSelectionMode(QListWidget.ExtendedSelection)
        merge_list_group.addWidget(self.list_widget)
        
        # 순서 및 제거 버튼
        control_buttons_layout = QVBoxLayout()
        self.up_button = QPushButton('▲ 위로')
        self.up_button.clicked.connect(self.move_up)
        control_buttons_layout.addWidget(self.up_button)
        self.down_button = QPushButton('▼ 아래로')
        self.down_button.clicked.connect(self.move_down)
        control_buttons_layout.addWidget(self.down_button)
        self.remove_button = QPushButton('X 제거')
        self.remove_button.clicked.connect(self.remove_file)
        control_buttons_layout.addWidget(self.remove_button)
        self.clear_button = QPushButton('전체 초기화')
        self.clear_button.clicked.connect(self.list_widget.clear)
        control_buttons_layout.addWidget(self.clear_button)
        control_buttons_layout.addStretch()
        
        merge_list_group.addLayout(control_buttons_layout)
        main_layout.addLayout(merge_list_group)

        # --- 3. 실행 및 설정 영역 ---
        settings_layout = QHBoxLayout()
        self.output_name_edit = QLineEdit("병합된_프레젠테이션.pptx")
        settings_layout.addWidget(QLabel("결과 파일 이름:"))
        settings_layout.addWidget(self.output_name_edit)
        
        self.save_path_label = QLabel(f"저장 경로: {os.path.basename(self.save_directory)}/... ")
        self.save_path_button = QPushButton('... 경로 선택')
        self.save_path_button.clicked.connect(self.select_save_path)
        settings_layout.addWidget(self.save_path_label)
        settings_layout.addWidget(self.save_path_button)
        main_layout.addLayout(settings_layout)

        self.merge_button = QPushButton('✅ 병합 실행')
        self.merge_button.setStyleSheet("font-size: 18px; padding: 10px; background-color: #4CAF50; color: white; border-radius: 5px;")
        self.merge_button.clicked.connect(self.execute_merge)
        main_layout.addWidget(self.merge_button)
        
        # --- 4. 상태 표시 영역 ---
        self.status_label = QLabel("상태: 파일을 추가하고 병합 순서를 지정해 주세요.")
        main_layout.addWidget(self.status_label)
        
        self.setLayout(main_layout)
        
        # 진행률 표시 대화 상자 초기화
        self.progress_dialog = QProgressDialog("PPTX 파일을 병합하는 중...", "취소", 0, 100, self)
        self.progress_dialog.setWindowTitle("병합 진행률")
        self.progress_dialog.setCancelButton(None) 
        self.progress_dialog.setWindowModality(Qt.WindowModal)
        self.progress_dialog.setAutoClose(False)
        self.progress_dialog.close()

    # --- 파일 탐색기로 파일 추가하는 기능 ---
    def open_file_dialog(self):
        filter_string = (
            "All Presentation Files (*.pptx *.ppt);;"
            "PPTX Files (*.pptx);;"                     
            "PPT Files (*.ppt);;"
            "All Files (*)"                             
        )
        
        file_names, _ = QFileDialog.getOpenFileNames(
            self, 
            '병합할 PPT 파일 선택', 
            '', 
            filter_string
        )
        
        if file_names:
            self.add_files_to_list(file_names)

    # --- 드래그 앤 드롭 기능 (오버라이드) ---
    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            urls = [url.path() for url in event.mimeData().urls()]
            is_valid_file = all(url.lower().endswith(('.pptx', '.ppt')) for url in urls)
            if is_valid_file:
                 event.accept()
            else:
                event.ignore()
        else:
            event.ignore()

    def dropEvent(self, event):
        file_paths = []
        for url in event.mimeData().urls():
            path = url.toLocalFile()
            if path.lower().endswith(('.pptx', '.ppt')):
                file_paths.append(path)
        
        if file_paths:
            self.add_files_to_list(file_paths)
            event.accept()
        else:
            event.ignore()

    # --- 파일 목록 관리 헬퍼 함수 ---
    def add_files_to_list(self, file_names):
        for file in file_names:
            if not self.list_widget.findItems(file, Qt.MatchExactly):
                self.list_widget.addItem(file)
        self.status_label.setText(f"상태: {self.list_widget.count()}개 파일이 추가되었습니다. 병합 준비 완료.")
        
    def move_up(self):
        current_row = self.list_widget.currentRow()
        if current_row > 0:
            item = self.list_widget.takeItem(current_row)
            self.list_widget.insertItem(current_row - 1, item)
            self.list_widget.setCurrentRow(current_row - 1)

    def move_down(self):
        current_row = self.list_widget.currentRow()
        if current_row < self.list_widget.count() - 1 and current_row != -1:
            item = self.list_widget.takeItem(current_row)
            self.list_widget.insertItem(current_row + 1, item)
            self.list_widget.setCurrentRow(current_row + 1)

    def remove_file(self):
        for item in self.list_widget.selectedItems():
            self.list_widget.takeItem(self.list_widget.row(item))
        
        self.status_label.setText(f"상태: {self.list_widget.count()}개 파일이 추가되었습니다. 병합 준비 완료.")

    def select_save_path(self):
        directory = QFileDialog.getExistingDirectory(self, '결과 파일 저장 경로 선택', self.save_directory)
        
        if directory:
            self.save_directory = directory
            self.save_path_label.setText(f"저장 경로: {os.path.basename(directory)}/... ")
        else:
            pass

    # --- 실제 병합 실행 함수 ---
    def execute_merge(self):
        file_paths = [self.list_widget.item(i).text() for i in range(self.list_widget.count())]
        output_filename = self.output_name_edit.text().strip()
        
        if not file_paths:
            QMessageBox.warning(self, "경고", "병합할 PPT 파일을 1개 이상 추가해야 합니다.")
            return

        if not output_filename:
            QMessageBox.warning(self, "경고", "결과 파일 이름을 입력해야 합니다.")
            return
            
        if not output_filename.lower().endswith(".pptx"):
            output_filename += ".pptx"

        if not self.save_directory:
            QMessageBox.warning(self, "경고", "저장 경로를 먼저 선택해 주세요.")
            return

        output_full_path = os.path.join(self.save_directory, output_filename)
        
        self.status_label.setText("상태: 병합 작업 시작 중...")
        self.merge_button.setEnabled(False)
        
        self.worker = MergerWorker(file_paths, output_full_path)
        self.worker.merge_finished.connect(self.on_merge_finished)
        self.worker.progress_update.connect(self.on_progress_update)
        self.worker.start()
        
        self.progress_dialog.show()
        
    # --- 워커 스레드 시그널 처리 ---
    def on_progress_update(self, current, total):
        self.progress_dialog.setMaximum(total)
        if current == 1 and total == 1:
            self.progress_dialog.setLabelText("PPT 파일을 PPTX로 변환 중...")
            self.progress_dialog.setMaximum(100)
            self.progress_dialog.setValue(50)
        else:
            self.progress_dialog.setValue(current)
            self.progress_dialog.setLabelText(f"슬라이드 복사 중: {current}/{total}")

    def on_merge_finished(self, success, message):
        self.merge_button.setEnabled(True)
        self.progress_dialog.close()
        
        if success:
            self.status_label.setText("상태: ✅ 병합 완료!")
            QMessageBox.information(self, "성공", message)
        else:
            self.status_label.setText("상태: ❌ 오류 발생!")
            QMessageBox.critical(self, "오류 발생", message)
            
        self.worker = None

# 프로그램 실행 부분
if __name__ == '__main__':
    app = QApplication(sys.argv)
    
    app.setStyleSheet("""
        QWidget { font-size: 10pt; }
        QPushButton { padding: 8px; border-radius: 5px; }
        QListWidget { border: 1px solid #ddd; padding: 5px; }
        QLabel#status_label { font-weight: bold; padding: 5px; }
    """)
    
    window = PptxMergerApp()
    window.show()
    sys.exit(app.exec_())
