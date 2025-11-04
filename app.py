# import sys
# from PyQt5.QtWidgets import (
#     QApplication, QWidget, QVBoxLayout, QHBoxLayout, 
#     QPushButton, QLineEdit, QLabel, QListWidget, 
#     QFileDialog, QSizePolicy
# )
# from PyQt5.QtCore import Qt
# # from pptx_merger import PptxMerger # (추후에 구현할 PPTX 병합 로직)

# class PptxMergerApp(QWidget):
#     def __init__(self):
#         super().__init__()
#         self.setWindowTitle('PPTX 순서 병합 프로그램')
#         self.setGeometry(100, 100, 600, 450)
#         self.setup_ui()
        
#         # 💡 드래그 앤 드롭 활성화
#         self.setAcceptDrops(True)

#     def setup_ui(self):
#         main_layout = QVBoxLayout()
        
#         # --- 1. 파일 추가 및 드래그 앤 드롭 영역 ---
#         file_select_layout = QHBoxLayout()
        
#         # "+ 파일 추가" 버튼 (파일 탐색기 열기)
#         self.select_files_button = QPushButton('+ 파일 추가 (탐색기)')
#         self.select_files_button.clicked.connect(self.open_file_dialog)
#         file_select_layout.addWidget(self.select_files_button)

#         # 드래그 앤 드롭 안내 레이블
#         drag_label = QLabel("여기에 파일 드래그 & 드롭 가능")
#         drag_label.setAlignment(Qt.AlignCenter)
#         drag_label.setStyleSheet("border: 2px dashed #ccc; padding: 10px;")
#         file_select_layout.addWidget(drag_label)
        
#         main_layout.addLayout(file_select_layout)
        
#         # --- 2. 병합 목록 영역 ---
#         merge_list_group = QHBoxLayout()
        
#         # 파일 목록 (List Widget)
#         self.list_widget = QListWidget()
#         merge_list_group.addWidget(self.list_widget)
        
#         # 순서 및 제거 버튼
#         control_buttons_layout = QVBoxLayout()
        
#         self.up_button = QPushButton('▲ 위로')
#         self.up_button.clicked.connect(self.move_up)
#         control_buttons_layout.addWidget(self.up_button)
        
#         self.down_button = QPushButton('▼ 아래로')
#         self.down_button.clicked.connect(self.move_down)
#         control_buttons_layout.addWidget(self.down_button)
        
#         self.remove_button = QPushButton('X 제거')
#         self.remove_button.clicked.connect(self.remove_file)
#         control_buttons_layout.addWidget(self.remove_button)
        
#         self.clear_button = QPushButton('전체 초기화')
#         self.clear_button.clicked.connect(self.list_widget.clear)
#         control_buttons_layout.addWidget(self.clear_button)
        
#         control_buttons_layout.addStretch() # 나머지 공간 채우기
        
#         merge_list_group.addLayout(control_buttons_layout)
#         main_layout.addLayout(merge_list_group)

#         # --- 3. 실행 및 설정 영역 ---
#         settings_layout = QHBoxLayout()
#         self.output_name_edit = QLineEdit("병합된_프레젠테이션.pptx")
#         settings_layout.addWidget(QLabel("결과 파일 이름:"))
#         settings_layout.addWidget(self.output_name_edit)
        
#         self.save_path_label = QLabel("저장 경로: (클릭 후 설정)")
#         self.save_path_button = QPushButton('... 경로 선택')
#         self.save_path_button.clicked.connect(self.select_save_path)
#         settings_layout.addWidget(self.save_path_label)
#         settings_layout.addWidget(self.save_path_button)
#         main_layout.addLayout(settings_layout)

#         self.merge_button = QPushButton('✅ 병합 실행')
#         self.merge_button.setStyleSheet("font-size: 18px; padding: 10px;")
#         # self.merge_button.clicked.connect(self.execute_merge) # 실제 병합 함수 연결 예정
#         main_layout.addWidget(self.merge_button)
        
#         # --- 4. 상태 표시 영역 ---
#         self.status_label = QLabel("상태: 파일을 추가해주세요.")
#         main_layout.addWidget(self.status_label)
        
#         self.setLayout(main_layout)

#     # --- 파일 탐색기로 파일 추가하는 기능 ---
#     # def open_file_dialog(self):
#     #     # PPTX 파일만 필터링하도록 설정
#     #     file_names, _ = QFileDialog.getOpenFileNames(
#     #         self, '병합할 PPTX 파일 선택', '', "PPTX Files (*.pptx);;PPT Files (*.ppt)"
#     #     )
#     #     if file_names:
#     #         self.add_files_to_list(file_names)


#     # --- 드래그 앤 드롭 기능 (오버라이드) ---
#     def dragEnterEvent(self, event):
#         if event.mimeData().hasUrls():
#             # URL (파일) 형식이고, 확장자가 pptx 또는 ppt 인 경우에만 드롭 허용
#             urls = [url.path() for url in event.mimeData().urls()]
#             is_pptx_file = all(url.lower().endswith(('.pptx', '.ppt')) for url in urls)
#             if is_pptx_file:
#                  event.accept()
#             else:
#                 event.ignore()
#         else:
#             event.ignore()

#     def dropEvent(self, event):
#         file_paths = []
#         for url in event.mimeData().urls():
#             # Mac에서는 경로 앞에 'file://' 같은 접두사가 붙으므로 경로만 추출
#             path = url.toLocalFile()
#             if path.lower().endswith(('.pptx', '.ppt')):
#                 file_paths.append(path)
        
#         if file_paths:
#             self.add_files_to_list(file_paths)
#             event.accept()
#         else:
#             event.ignore()

#     # --- 파일 목록 관리 헬퍼 함수 ---
#     def add_files_to_list(self, file_names):
#         for file in file_names:
#             # 중복 추가 방지 로직 (선택 사항)
#             # if self.list_widget.findItems(file, Qt.MatchExactly):
#             #     continue
#             self.list_widget.addItem(file)
#         self.status_label.setText(f"상태: {self.list_widget.count()}개 파일이 추가되었습니다.")
        
#     def move_up(self):
#         # 현재 선택된 아이템을 가져옵니다.
#         current_row = self.list_widget.currentRow()
        
#         # 첫 번째 항목이 아니어야만 위로 이동할 수 있습니다.
#         if current_row > 0:
#             # 선택된 아이템을 제거하고 (떼어내고), 바로 위 인덱스에 다시 삽입합니다.
#             item = self.list_widget.takeItem(current_row)
#             self.list_widget.insertItem(current_row - 1, item)
            
#             # 이동한 아이템을 다시 선택 상태로 만듭니다.
#             self.list_widget.setCurrentRow(current_row - 1)

#     def move_down(self):
#         current_row = self.list_widget.currentRow()
#         # 마지막 항목이 아니어야만 아래로 이동할 수 있습니다.
#         if current_row < self.list_widget.count() - 1 and current_row != -1:
#             # 선택된 아이템을 제거하고, 바로 아래 인덱스에 다시 삽입합니다.
#             item = self.list_widget.takeItem(current_row)
#             self.list_widget.insertItem(current_row + 1, item)
            
#             # 이동한 아이템을 다시 선택 상태로 만듭니다.
#             self.list_widget.setCurrentRow(current_row + 1)

#     def remove_file(self):
#         # 현재 선택된 모든 아이템을 가져옵니다.
#         for item in self.list_widget.selectedItems():
#             # 아이템이 속한 행(row)을 찾아서 제거합니다.
#             self.list_widget.takeItem(self.list_widget.row(item))
        
#         # 상태 표시 업데이트
#         self.status_label.setText(f"상태: {self.list_widget.count()}개 파일이 추가되었습니다.")

#     def select_save_path(self):
#         # 폴더 탐색기를 열어 저장할 경로(디렉토리)를 선택합니다.
#         directory = QFileDialog.getExistingDirectory(self, '결과 파일 저장 경로 선택', '')
        
#         if directory:
#             self.save_directory = directory
#             self.save_path_label.setText(f"저장 경로: {directory}")
#         else:
#             self.save_directory = None
#             self.save_path_label.setText("저장 경로: (클릭 후 설정)")

# # 프로그램 실행 부분
# if __name__ == '__main__':
#     app = QApplication(sys.argv)
#     window = PptxMergerApp()
#     window.show()
#     sys.exit(app.exec_())
import sys
import os
import traceback
from PyQt5.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QHBoxLayout, 
    QPushButton, QLineEdit, QLabel, QListWidget, 
    QFileDialog, QMessageBox, QSizePolicy, QProgressDialog
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, QCoreApplication
from pptx import Presentation
from pptx.util import Inches

# PPTX 병합 작업을 별도의 스레드에서 처리하기 위한 워커 클래스
class MergerWorker(QThread):
    # 작업 진행 상황을 GUI에 알리기 위한 시그널
    # 1. merge_finished: (성공 여부, 메시지)
    merge_finished = pyqtSignal(bool, str)
    # 2. progress_update: (현재 슬라이드 번호, 전체 슬라이드 개수)
    progress_update = pyqtSignal(int, int)

    def __init__(self, file_paths, output_path, parent=None):
        super().__init__(parent)
        self.file_paths = file_paths
        self.output_path = output_path

    def run(self):
        # 작업 시작 시그널을 보내거나, 상태를 "병합 중"으로 업데이트
        
        # 1. 파일 목록이 비어있는지 확인
        if not self.file_paths:
            self.merge_finished.emit(False, "오류: 병합할 PPT 파일이 없습니다.")
            return

        try:
            # 2. 첫 번째 파일을 마스터 프레젠테이션으로 로드
            # python-pptx는 .ppt 형식은 지원하지 않으므로, 사용자에게 .pptx 파일만 사용하도록 경고하거나,
            # .ppt 파일은 미리 .pptx로 변환해야 합니다. (여기서는 .pptx만 처리한다고 가정합니다)
            master_pptx = Presentation(self.file_paths[0])
            total_slides_processed = master_pptx.slides.count
            
            # 전체 슬라이드 개수 계산 (진행률 표시용)
            total_slides_count = sum(Presentation(path).slides.count for path in self.file_paths)
            
            # 진행률 업데이트
            self.progress_update.emit(total_slides_processed, total_slides_count)
            
            # 3. 나머지 파일들을 순회하며 슬라이드 복사
            for i, path in enumerate(self.file_paths[1:]):
                source_pptx = Presentation(path)
                
                # 마스터 슬라이드 정보 (레이아웃 복사용)
                # python-pptx에서 슬라이드 복사는 레이아웃을 참조하므로, 마스터 슬라이드를 참조하여 추가합니다.
                slide_layout_map = {layout.name: layout for layout in master_pptx.slide_layouts}

                for slide in source_pptx.slides:
                    # 마스터 PPTX에 존재하는 이름의 레이아웃을 사용
                    source_layout_name = slide.slide_layout.name
                    target_layout = slide_layout_map.get(source_layout_name, master_pptx.slide_layouts[6]) # 6: Blank Layout
                    
                    # 새로운 슬라이드를 마스터 프레젠테이션에 추가
                    new_slide = master_pptx.slides.add_slide(target_layout)
                    
                    # 콘텐츠 복사 (텍스트, 이미지 등)
                    # 텍스트 복사: 셰이프를 순회하며 텍스트 복사 (완벽하지 않을 수 있음)
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text_frame = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height).text_frame
                            text_frame.text = shape.text
                        elif shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE (이미지)
                            # 그림 파일은 복사 로직이 복잡하므로 간단히 건너뛰거나, 
                            # 임시 파일로 저장 후 다시 삽입해야 합니다. 
                            # 여기서는 간단히 건너뛰고 텍스트/표/차트는 기본 복사합니다.
                            pass

                    # Note: python-pptx의 슬라이드 복사는 매우 제한적입니다. 
                    # 복잡한 차트, 표, 이미지 등을 완벽하게 복사하려면 더 많은 코드가 필요합니다.
                    # 이 코드는 텍스트 기반 슬라이드에 대해 기본적으로 작동합니다.
                    
                    total_slides_processed += 1
                    # 진행률 업데이트
                    self.progress_update.emit(total_slides_processed, total_slides_count)
                    
                    # GUI 업데이트를 위한 강제 이벤트 처리
                    QCoreApplication.processEvents()

            # 4. 결과 파일 저장
            master_pptx.save(self.output_path)
            
            # 5. 성공 시그널
            self.merge_finished.emit(True, f"✅ PPTX 병합 완료! \n\n저장 위치: {self.output_path}")

        except Exception as e:
            error_message = f"PPTX 병합 중 오류가 발생했습니다.\n\n오류: {e}\n\n상세:\n{traceback.format_exc()}"
            self.merge_finished.emit(False, error_message)

class PptxMergerApp(QWidget):
    def __init__(self):
        super().__init__()
        # Pypptx는 .ppt를 지원하지 않으므로 제목을 명확히 합니다.
        self.setWindowTitle('PPTX 순서 병합 프로그램 (pptx 전용)') 
        self.setGeometry(100, 100, 650, 480)
        self.save_directory = os.path.expanduser("~") # 초기 저장 경로 설정
        self.setup_ui()
        
        # 💡 드래그 앤 드롭 활성화
        self.setAcceptDrops(True)
        
        self.worker = None # Worker 스레드 객체

    def setup_ui(self):
        main_layout = QVBoxLayout()
        
        # --- 1. 파일 추가 및 드래그 앤 드롭 영역 ---
        file_select_layout = QHBoxLayout()
        
        # "+ 파일 추가" 버튼 (파일 탐색기 열기)
        self.select_files_button = QPushButton('+ 파일 추가 (탐색기)')
        self.select_files_button.clicked.connect(self.open_file_dialog)
        file_select_layout.addWidget(self.select_files_button)

        # 드래그 앤 드롭 안내 레이블
        drag_label = QLabel("여기에 .PPTX 파일을 드래그 & 드롭 가능")
        drag_label.setAlignment(Qt.AlignCenter)
        drag_label.setStyleSheet("border: 2px dashed #ccc; padding: 10px; color: #555; background-color: #f9f9f9; border-radius: 8px;")
        file_select_layout.addWidget(drag_label)
        
        main_layout.addLayout(file_select_layout)
        
        # --- 2. 병합 목록 영역 ---
        merge_list_group = QHBoxLayout()
        
        # 파일 목록 (List Widget) - 선택 모드 설정
        self.list_widget = QListWidget()
        self.list_widget.setSelectionMode(QListWidget.ExtendedSelection) # 다중 선택 가능
        merge_list_group.addWidget(self.list_widget)
        
        # 순서 및 제거 버튼
        control_buttons_layout = QVBoxLayout()
        
        self.up_button = QPushButton('▲ 위로')
        self.up_button.clicked.connect(self.move_up)
        control_buttons_layout.addWidget(self.up_button)
        
        self.down_button = QPushButton('▼ 아래로')
        self.down_button.clicked.connect(self.move_down)
        control_buttons_layout.addWidget(self.down_button)
        
        self.remove_button = QPushButton('X 제거')
        self.remove_button.clicked.connect(self.remove_file)
        control_buttons_layout.addWidget(self.remove_button)
        
        self.clear_button = QPushButton('전체 초기화')
        self.clear_button.clicked.connect(self.list_widget.clear)
        control_buttons_layout.addWidget(self.clear_button)
        
        control_buttons_layout.addStretch() # 나머지 공간 채우기
        
        merge_list_group.addLayout(control_buttons_layout)
        main_layout.addLayout(merge_list_group)

        # --- 3. 실행 및 설정 영역 ---
        settings_layout = QHBoxLayout()
        self.output_name_edit = QLineEdit("병합된_프레젠테이션.pptx")
        settings_layout.addWidget(QLabel("결과 파일 이름:"))
        settings_layout.addWidget(self.output_name_edit)
        
        self.save_path_label = QLabel(f"저장 경로: {self.save_directory}")
        self.save_path_button = QPushButton('... 경로 선택')
        self.save_path_button.clicked.connect(self.select_save_path)
        settings_layout.addWidget(self.save_path_label)
        settings_layout.addWidget(self.save_path_button)
        main_layout.addLayout(settings_layout)

        self.merge_button = QPushButton('✅ 병합 실행')
        self.merge_button.setStyleSheet("font-size: 18px; padding: 10px; background-color: #4CAF50; color: white; border-radius: 5px;")
        self.merge_button.clicked.connect(self.execute_merge) # 실제 병합 함수 연결
        main_layout.addWidget(self.merge_button)
        
        # --- 4. 상태 표시 영역 ---
        self.status_label = QLabel("상태: 파일을 추가하고 병합 순서를 지정해 주세요.")
        main_layout.addWidget(self.status_label)
        
        self.setLayout(main_layout)
        
        # 진행률 표시 대화 상자 초기화
        self.progress_dialog = QProgressDialog("PPTX 파일을 병합하는 중...", "취소", 0, 100, self)
        self.progress_dialog.setWindowTitle("병합 진행률")
        self.progress_dialog.setCancelButton(None) # 취소 버튼 비활성화
        self.progress_dialog.setWindowModality(Qt.WindowModal)
        self.progress_dialog.setAutoClose(False)
        self.progress_dialog.close()

    # --- 파일 탐색기로 파일 추가하는 기능 (수정 완료) ---
    def open_file_dialog(self):
        # PPTX 파일만 지원하며, QFileDialog에서 필터가 올바르게 작동하도록 수정했습니다.
        filter_string = (
            "PPTX Files (*.pptx);;"                     
            "All Files (*)"                             
        )
        
        file_names, _ = QFileDialog.getOpenFileNames(
            self, 
            '병합할 PPTX 파일 선택', 
            '', 
            filter_string
        )
        
        # 파일이 선택된 경우 목록에 추가
        if file_names:
            self.add_files_to_list(file_names)

    # --- 드래그 앤 드롭 기능 (오버라이드) ---
    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            # URL (파일) 형식이고, 확장자가 pptx 인 경우에만 드롭 허용
            urls = [url.path() for url in event.mimeData().urls()]
            is_pptx_file = all(url.lower().endswith('.pptx') for url in urls)
            if is_pptx_file:
                 event.accept()
            else:
                event.ignore()
        else:
            event.ignore()

    def dropEvent(self, event):
        file_paths = []
        for url in event.mimeData().urls():
            # 경로만 추출 및 pptx 필터링
            path = url.toLocalFile()
            if path.lower().endswith('.pptx'):
                file_paths.append(path)
        
        if file_paths:
            self.add_files_to_list(file_paths)
            event.accept()
        else:
            event.ignore()

    # --- 파일 목록 관리 헬퍼 함수 ---
    def add_files_to_list(self, file_names):
        for file in file_names:
            # 중복 추가 방지
            if not self.list_widget.findItems(file, Qt.MatchExactly):
                self.list_widget.addItem(file)
        self.status_label.setText(f"상태: {self.list_widget.count()}개 파일이 추가되었습니다. 병합 준비 완료.")
        
    def move_up(self):
        current_row = self.list_widget.currentRow()
        if current_row > 0:
            item = self.list_widget.takeItem(current_row)
            self.list_widget.insertItem(current_row - 1, item)
            self.list_widget.setCurrentRow(current_row - 1)

    def move_down(self):
        current_row = self.list_widget.currentRow()
        if current_row < self.list_widget.count() - 1 and current_row != -1:
            item = self.list_widget.takeItem(current_row)
            self.list_widget.insertItem(current_row + 1, item)
            self.list_widget.setCurrentRow(current_row + 1)

    def remove_file(self):
        for item in self.list_widget.selectedItems():
            self.list_widget.takeItem(self.list_widget.row(item))
        
        self.status_label.setText(f"상태: {self.list_widget.count()}개 파일이 추가되었습니다. 병합 준비 완료.")

    def select_save_path(self):
        # 폴더 탐색기를 열어 저장할 경로(디렉토리)를 선택합니다.
        directory = QFileDialog.getExistingDirectory(self, '결과 파일 저장 경로 선택', self.save_directory)
        
        if directory:
            self.save_directory = directory
            self.save_path_label.setText(f"저장 경로: {os.path.basename(directory)}/... ")
        else:
            # 사용자가 취소했을 경우 현재 경로를 유지합니다.
            pass

    # --- 실제 병합 실행 함수 ---
    def execute_merge(self):
        # 1. 입력 검증
        file_paths = [self.list_widget.item(i).text() for i in range(self.list_widget.count())]
        output_filename = self.output_name_edit.text().strip()
        
        if not file_paths:
            QMessageBox.warning(self, "경고", "병합할 PPTX 파일을 1개 이상 추가해야 합니다.")
            return

        if not output_filename:
            QMessageBox.warning(self, "경고", "결과 파일 이름을 입력해야 합니다.")
            return
            
        # .pptx 확장자 자동 추가
        if not output_filename.lower().endswith(".pptx"):
            output_filename += ".pptx"

        if not self.save_directory:
            QMessageBox.warning(self, "경고", "저장 경로를 먼저 선택해 주세요.")
            return

        output_full_path = os.path.join(self.save_directory, output_filename)
        
        # 2. 작업 시작 및 GUI 비활성화
        self.status_label.setText("상태: 병합 작업 시작 중...")
        self.merge_button.setEnabled(False)
        
        # 3. 워커 스레드 생성 및 실행
        self.worker = MergerWorker(file_paths, output_full_path)
        self.worker.merge_finished.connect(self.on_merge_finished)
        self.worker.progress_update.connect(self.on_progress_update)
        self.worker.start()
        
        # 4. 진행률 다이얼로그 표시
        self.progress_dialog.show()
        
    # --- 워커 스레드 시그널 처리 ---
    def on_progress_update(self, current, total):
        self.progress_dialog.setMaximum(total)
        self.progress_dialog.setValue(current)
        self.progress_dialog.setLabelText(f"슬라이드 복사 중: {current}/{total}")

    def on_merge_finished(self, success, message):
        # 1. GUI 상태 복구
        self.merge_button.setEnabled(True)
        self.progress_dialog.close()
        
        # 2. 결과 처리
        if success:
            self.status_label.setText("상태: ✅ 병합 완료! (자세한 내용은 메시지 확인)")
            QMessageBox.information(self, "성공", message)
        else:
            self.status_label.setText("상태: ❌ 오류 발생! (자세한 내용은 메시지 확인)")
            QMessageBox.critical(self, "오류 발생", message)
            
        # 3. 워커 객체 정리
        self.worker = None

# 프로그램 실행 부분
if __name__ == '__main__':
    # QCoreApplication.setAttribute(Qt.AA_EnableHighDpiScaling) # DPI 스케일링 활성화 (옵션)
    app = QApplication(sys.argv)
    
    # 폰트 및 스타일링 (선택 사항)
    app.setStyleSheet("""
        QWidget { font-size: 10pt; }
        QPushButton { padding: 8px; border-radius: 5px; }
        QListWidget { border: 1px solid #ddd; padding: 5px; }
        QLabel#status_label { font-weight: bold; padding: 5px; }
    """)
    
    window = PptxMergerApp()
    window.show()
    sys.exit(app.exec_())
