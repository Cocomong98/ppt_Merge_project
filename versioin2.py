import sys
import os
import traceback
from PyQt5.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QHBoxLayout, 
    QPushButton, QLineEdit, QLabel, QListWidget, 
    QFileDialog, QMessageBox, QProgressDialog
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, QCoreApplication
from pptx import Presentation

# win32com ê´€ë ¨ ëª¨ë“ˆì„ Windows í™˜ê²½ì—ì„œë§Œ importí•˜ë„ë¡ ì²˜ë¦¬
try:
    import pythoncom
    import win32com.client
    WINDOWS_ENV = True
except ImportError:
    # Mac ê°œë°œ í™˜ê²½ì´ê±°ë‚˜ pywin32ê°€ ì„¤ì¹˜ë˜ì§€ ì•Šì€ ê²½ìš°
    WINDOWS_ENV = False 

# PPTX ë³‘í•© ìž‘ì—…ì„ ë³„ë„ì˜ ìŠ¤ë ˆë“œì—ì„œ ì²˜ë¦¬í•˜ê¸° ìœ„í•œ ì›Œì»¤ í´ëž˜ìŠ¤
class MergerWorker(QThread):
    merge_finished = pyqtSignal(bool, str) # (ì„±ê³µ ì—¬ë¶€, ë©”ì‹œì§€)
    progress_update = pyqtSignal(int, int) # (í˜„ìž¬ ìŠ¬ë¼ì´ë“œ ë²ˆí˜¸, ì „ì²´ ìŠ¬ë¼ì´ë“œ ê°œìˆ˜)

    def __init__(self, file_paths, output_path, parent=None):
        super().__init__(parent)
        self.file_paths = file_paths
        self.output_path = output_path
        self.temp_files = [] # ë³€í™˜ëœ ìž„ì‹œ íŒŒì¼ ëª©ë¡

    def __del__(self):
        # ìŠ¤ë ˆë“œ ì¢…ë£Œ ì‹œ ìž„ì‹œ íŒŒì¼ ì •ë¦¬
        for temp_file in self.temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception:
                pass

    # .ppt íŒŒì¼ì„ .pptxë¡œ ë³€í™˜í•˜ëŠ” í•¨ìˆ˜ (Windows + MS PowerPoint í•„ìš”)
    def convert_ppt_to_pptx(self, ppt_path):
        if not WINDOWS_ENV:
            raise EnvironmentError("`.ppt` íŒŒì¼ ë³€í™˜ì€ Windows í™˜ê²½ê³¼ `pywin32` ë¼ì´ë¸ŒëŸ¬ë¦¬, ê·¸ë¦¬ê³  MS PowerPointê°€ í•„ìš”í•©ë‹ˆë‹¤.")
        
        # COM ê°ì²´ ì‚¬ìš© ì „ ì´ˆê¸°í™”
        pythoncom.CoInitialize() 
        
        temp_dir = os.path.dirname(ppt_path)
        # ìž„ì‹œ íŒŒì¼ ì´ë¦„ì€ ì›ë³¸ê³¼ ë™ì¼í•œ ê²½ë¡œì— ìƒì„±
        temp_pptx_path = os.path.join(temp_dir, f"~temp_converted_{os.path.basename(ppt_path)[:-4]}.pptx")
        
        powerpoint = None
        try:
            # PowerPoint ê°ì²´ ìƒì„±
            powerpoint = win32com.client.Dispatch("Powerpoint.Application")
            
            # ðŸ’¡ ì°½ì„ ëª…ì‹œì ìœ¼ë¡œ ë³´ì´ê²Œ í•˜ì—¬ ê¶Œí•œ ë¬¸ì œ íšŒí”¼
            powerpoint.Visible = True  
            
            # íŒŒì¼ì„ ì—´ ë•Œ ë°œìƒí•˜ëŠ” ê²½ê³ ë¥¼ ë¬´ì‹œí•˜ê³  ì•ˆì •ì„± í–¥ìƒ
            presentation = powerpoint.Presentations.Open(
                ppt_path, 
                ReadOnly=True,
                Untitled=False,      
                WithAddins=False     
            )
            
            presentation.SaveAs(temp_pptx_path, 24) # 24ëŠ” ppSaveAsPresentation (pptx)
            presentation.Close()
            
            self.temp_files.append(temp_pptx_path)
            return temp_pptx_path
        except Exception as e:
            # ì˜¤ë¥˜ ë°œìƒ ì‹œ PowerPoint ê°ì²´ë¥¼ í™•ì‹¤ížˆ ì¢…ë£Œ
            if powerpoint:
                try: powerpoint.Quit()
                except: pass
            
            # ì˜¤ë¥˜ ë©”ì‹œì§€ì— ì •í™•í•œ ì˜ˆì™¸ ì½”ë“œë¥¼ í¬í•¨í•˜ì—¬ ì‚¬ìš©ìžì—ê²Œ ì „ë‹¬
            error_details = f"PPT íŒŒì¼ ë³€í™˜ ì‹¤íŒ¨ (MS PowerPoint ì„¤ì¹˜ ë° ê¶Œí•œ í™•ì¸ í•„ìš”): {e}"
            raise Exception(error_details)
        finally:
            # ìŠ¤ë ˆë“œ ì¢…ë£Œ ì‹œ COM Uninitialization
            pythoncom.CoUninitialize()

    def run(self):
        if not self.file_paths:
            self.merge_finished.emit(False, "ì˜¤ë¥˜: ë³‘í•©í•  PPT íŒŒì¼ì´ ì—†ìŠµë‹ˆë‹¤.")
            return
            
        process_paths = []
        try:
            # 1. íŒŒì¼ ëª©ë¡ ìˆœíšŒ ë° PPTXë¡œ ë³€í™˜
            for path in self.file_paths:
                if path.lower().endswith('.ppt'):
                    # win32comì„ ì‚¬ìš©í•˜ë ¤ë©´ Windows í™˜ê²½ì´ì–´ì•¼ í•¨ì„ ë‹¤ì‹œ í™•ì¸
                    if not WINDOWS_ENV:
                         raise EnvironmentError("ì´ PCëŠ” Windows í™˜ê²½ì´ ì•„ë‹™ë‹ˆë‹¤. `.ppt` ë³€í™˜ì€ Windowsì—ì„œë§Œ ê°€ëŠ¥í•©ë‹ˆë‹¤.")

                    self.progress_update.emit(1, 1) # ë³€í™˜ ë‹¨ê³„ í‘œì‹œ
                    converted_path = self.convert_ppt_to_pptx(path)
                    process_paths.append(converted_path)
                elif path.lower().endswith('.pptx'):
                    process_paths.append(path)
                
            if not process_paths:
                self.merge_finished.emit(False, "ì˜¤ë¥˜: ë³‘í•©í•  íŒŒì¼ì´ ìœ íš¨í•˜ì§€ ì•ŠìŠµë‹ˆë‹¤.")
                return

            # 2. ë³‘í•© ë¡œì§ ì‹œìž‘ ë° ìŠ¬ë¼ì´ë“œ ì¹´ìš´íŠ¸ ìµœì í™”
            master_pptx = Presentation(process_paths[0])
            
            total_slides_processed = len(master_pptx.slides) 
            total_slides_count = total_slides_processed
            
            # ë‚˜ë¨¸ì§€ íŒŒì¼ë“¤ì˜ ìŠ¬ë¼ì´ë“œ ìˆ˜ë§Œ í•©ì‚°í•˜ì—¬ ì „ì²´ ì¹´ìš´íŠ¸ ê³„ì‚°
            for path in process_paths[1:]:
                # íŒŒì¼ì„ ë¡œë“œí•˜ì—¬ ìŠ¬ë¼ì´ë“œ ê°œìˆ˜ë§Œ ì–»ì–´ì˜´
                total_slides_count += len(Presentation(path).slides) 

            self.progress_update.emit(total_slides_processed, total_slides_count)
            
            # 3. ë‚˜ë¨¸ì§€ íŒŒì¼ë“¤ì„ ìˆœíšŒí•˜ë©° ìŠ¬ë¼ì´ë“œ ë³µì‚¬
            for path in process_paths[1:]:
                source_pptx = Presentation(path)
                slide_layout_map = {layout.name: layout for layout in master_pptx.slide_layouts}

                for slide in source_pptx.slides:
                    source_layout_name = slide.slide_layout.name
                    target_layout = slide_layout_map.get(source_layout_name, master_pptx.slide_layouts[6])
                    
                    new_slide = master_pptx.slides.add_slide(target_layout)
                    
                    # ì½˜í…ì¸  ë³µì‚¬ (í…ìŠ¤íŠ¸ë§Œ ë³µì‚¬, ì´ë¯¸ì§€/ì°¨íŠ¸ëŠ” ìƒëžµ)
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            try:
                                text_frame = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height).text_frame
                                text_frame.text = shape.text
                            except Exception:
                                # í¬ê¸°/ìœ„ì¹˜ ì˜¤ë¥˜ ì‹œ ì•ˆì „í•˜ê²Œ í…ìŠ¤íŠ¸ë§Œ ë³µì‚¬
                                new_slide.shapes.add_textbox(0, 0, 1, 1).text_frame.text = shape.text
                        elif shape.shape_type == 13: 
                            pass # ì´ë¯¸ì§€/ì°¨íŠ¸ ìƒëžµ

                    total_slides_processed += 1
                    self.progress_update.emit(total_slides_processed, total_slides_count)
                    QCoreApplication.processEvents()

            # 4. ê²°ê³¼ íŒŒì¼ ì €ìž¥
            master_pptx.save(self.output_path)
            
            self.merge_finished.emit(True, f"âœ… PPTX ë³‘í•© ì™„ë£Œ!\n\nì €ìž¥ ìœ„ì¹˜: {self.output_path}")

        except Exception as e:
            error_message = f"PPTX ë³‘í•© ì¤‘ ì˜¤ë¥˜ê°€ ë°œìƒí–ˆìŠµë‹ˆë‹¤.\n\nì˜¤ë¥˜: {e}\n\nìƒì„¸:\n{traceback.format_exc()}"
            self.merge_finished.emit(False, error_message)
        finally:
            # ìž„ì‹œ íŒŒì¼ ì •ë¦¬
            for temp_file in self.temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except Exception:
                    pass

class PptxMergerApp(QWidget):
    def __init__(self):
        super().__init__()
        # win32com ì‚¬ìš© ì—¬ë¶€ì— ë”°ë¼ íƒ€ì´í‹€ ëª…ì‹œ
        title_suffix = " (MS PowerPoint ì œì–´)" if WINDOWS_ENV else " (pywin32 ì„¤ì¹˜ í•„ìš”)"
        self.setWindowTitle('PPTX/PPT ìˆœì„œ ë³‘í•© í”„ë¡œê·¸ëž¨' + title_suffix) 
        self.setGeometry(100, 100, 650, 480)
        self.save_directory = os.path.expanduser("~") 
        self.setup_ui()
        
        self.setAcceptDrops(True)
        self.worker = None

    def setup_ui(self):
        main_layout = QVBoxLayout()
        
        # --- 1. íŒŒì¼ ì¶”ê°€ ë° ë“œëž˜ê·¸ ì•¤ ë“œë¡­ ì˜ì—­ ---
        file_select_layout = QHBoxLayout()
        
        self.select_files_button = QPushButton('+ íŒŒì¼ ì¶”ê°€ (íƒìƒ‰ê¸°)')
        self.select_files_button.clicked.connect(self.open_file_dialog)
        file_select_layout.addWidget(self.select_files_button)

        drag_info = "ì—¬ê¸°ì— **.PPTX ë˜ëŠ” .PPT** íŒŒì¼ì„ ë“œëž˜ê·¸ & ë“œë¡­ ê°€ëŠ¥"
        if WINDOWS_ENV:
            drag_info += "\n(âš ï¸ .PPT ë³€í™˜: MS PowerPoint ì„¤ì¹˜ ë° ê´€ë¦¬ìž ê¶Œí•œ í•„ìˆ˜)"
        
        drag_label = QLabel(drag_info)
        drag_label.setAlignment(Qt.AlignCenter)
        drag_label.setStyleSheet("border: 2px dashed #ccc; padding: 10px; color: #1a1a1a; background-color: #e6f7ff; border-radius: 8px; font-weight: bold;")
        file_select_layout.addWidget(drag_label)
        
        main_layout.addLayout(file_select_layout)
        
        # --- 2. ë³‘í•© ëª©ë¡ ì˜ì—­ ---
        merge_list_group = QHBoxLayout()
        
        self.list_widget = QListWidget()
        self.list_widget.setSelectionMode(QListWidget.ExtendedSelection)
        merge_list_group.addWidget(self.list_widget)
        
        # ìˆœì„œ ë° ì œê±° ë²„íŠ¼
        control_buttons_layout = QVBoxLayout()
        self.up_button = QPushButton('â–² ìœ„ë¡œ')
        self.up_button.clicked.connect(self.move_up)
        control_buttons_layout.addWidget(self.up_button)
        self.down_button = QPushButton('â–¼ ì•„ëž˜ë¡œ')
        self.down_button.clicked.connect(self.move_down)
        control_buttons_layout.addWidget(self.down_button)
        self.remove_button = QPushButton('X ì œê±°')
        self.remove_button.clicked.connect(self.remove_file)
        control_buttons_layout.addWidget(self.remove_button)
        self.clear_button = QPushButton('ì „ì²´ ì´ˆê¸°í™”')
        self.clear_button.clicked.connect(self.list_widget.clear)
        control_buttons_layout.addWidget(self.clear_button)
        control_buttons_layout.addStretch()
        
        merge_list_group.addLayout(control_buttons_layout)
        main_layout.addLayout(merge_list_group)

        # --- 3. ì‹¤í–‰ ë° ì„¤ì • ì˜ì—­ ---
        settings_layout = QHBoxLayout()
        self.output_name_edit = QLineEdit("ë³‘í•©ëœ_í”„ë ˆì  í…Œì´ì…˜.pptx")
        settings_layout.addWidget(QLabel("ê²°ê³¼ íŒŒì¼ ì´ë¦„:"))
        settings_layout.addWidget(self.output_name_edit)
        
        self.save_path_label = QLabel(f"ì €ìž¥ ê²½ë¡œ: {os.path.basename(self.save_directory)}/... ")
        self.save_path_button = QPushButton('... ê²½ë¡œ ì„ íƒ')
        self.save_path_button.clicked.connect(self.select_save_path)
        settings_layout.addWidget(self.save_path_label)
        settings_layout.addWidget(self.save_path_button)
        main_layout.addLayout(settings_layout)

        self.merge_button = QPushButton('âœ… ë³‘í•© ì‹¤í–‰')
        self.merge_button.setStyleSheet("font-size: 18px; padding: 10px; background-color: #4CAF50; color: white; border-radius: 5px;")
        self.merge_button.clicked.connect(self.execute_merge)
        main_layout.addWidget(self.merge_button)
        
        # --- 4. ìƒíƒœ í‘œì‹œ ì˜ì—­ ---
        self.status_label = QLabel("ìƒíƒœ: íŒŒì¼ì„ ì¶”ê°€í•˜ê³  ë³‘í•© ìˆœì„œë¥¼ ì§€ì •í•´ ì£¼ì„¸ìš”.")
        main_layout.addWidget(self.status_label)
        
        self.setLayout(main_layout)
        
        # ì§„í–‰ë¥  í‘œì‹œ ëŒ€í™” ìƒìž ì´ˆê¸°í™”
        self.progress_dialog = QProgressDialog("PPTX íŒŒì¼ì„ ë³‘í•©í•˜ëŠ” ì¤‘...", "ì·¨ì†Œ", 0, 100, self)
        self.progress_dialog.setWindowTitle("ë³‘í•© ì§„í–‰ë¥ ")
        self.progress_dialog.setCancelButton(None) 
        self.progress_dialog.setWindowModality(Qt.WindowModal)
        self.progress_dialog.setAutoClose(False)
        self.progress_dialog.close()

    # --- íŒŒì¼ íƒìƒ‰ê¸°ë¡œ íŒŒì¼ ì¶”ê°€í•˜ëŠ” ê¸°ëŠ¥ ---
    def open_file_dialog(self):
        filter_string = (
            "All Presentation Files (*.pptx *.ppt);;"
            "PPTX Files (*.pptx);;"                     
            "PPT Files (*.ppt);;"
            "All Files (*)"                             
        )
        
        file_names, _ = QFileDialog.getOpenFileNames(
            self, 
            'ë³‘í•©í•  PPT íŒŒì¼ ì„ íƒ', 
            '', 
            filter_string
        )
        
        if file_names:
            self.add_files_to_list(file_names)

    # --- ë“œëž˜ê·¸ ì•¤ ë“œë¡­ ê¸°ëŠ¥ (ì˜¤ë²„ë¼ì´ë“œ) ---
    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            urls = [url.path() for url in event.mimeData().urls()]
            # .pptxì™€ .ppt ëª¨ë‘ í—ˆìš©í•©ë‹ˆë‹¤.
            is_valid_file = all(url.lower().endswith(('.pptx', '.ppt')) for url in urls)
            if is_valid_file:
                 event.accept()
            else:
                event.ignore()
        else:
            event.ignore()

    def dropEvent(self, event):
        file_paths = []
        for url in event.mimeData().urls():
            path = url.toLocalFile()
            if path.lower().endswith(('.pptx', '.ppt')):
                file_paths.append(path)
        
        if file_paths:
            self.add_files_to_list(file_paths)
            event.accept()
        else:
            event.ignore()

    # --- íŒŒì¼ ëª©ë¡ ê´€ë¦¬ í—¬í¼ í•¨ìˆ˜ (ë‚˜ë¨¸ì§€ ë¶€ë¶„ì€ ë™ì¼) ---
    def add_files_to_list(self, file_names):
        for file in file_names:
            if not self.list_widget.findItems(file, Qt.MatchExactly):
                self.list_widget.addItem(file)
        self.status_label.setText(f"ìƒíƒœ: {self.list_widget.count()}ê°œ íŒŒì¼ì´ ì¶”ê°€ë˜ì—ˆìŠµë‹ˆë‹¤. ë³‘í•© ì¤€ë¹„ ì™„ë£Œ.")
        
    def move_up(self):
        current_row = self.list_widget.currentRow()
        if current_row > 0:
            item = self.list_widget.takeItem(current_row)
            self.list_widget.insertItem(current_row - 1, item)
            self.list_widget.setCurrentRow(current_row - 1)

    def move_down(self):
        current_row = self.list_widget.currentRow()
        if current_row < self.list_widget.count() - 1 and current_row != -1:
            item = self.list_widget.takeItem(current_row)
            self.list_widget.insertItem(current_row + 1, item)
            self.list_widget.setCurrentRow(current_row + 1)

    def remove_file(self):
        for item in self.list_widget.selectedItems():
            self.list_widget.takeItem(self.list_widget.row(item))
        
        self.status_label.setText(f"ìƒíƒœ: {self.list_widget.count()}ê°œ íŒŒì¼ì´ ì¶”ê°€ë˜ì—ˆìŠµë‹ˆë‹¤. ë³‘í•© ì¤€ë¹„ ì™„ë£Œ.")

    def select_save_path(self):
        directory = QFileDialog.getExistingDirectory(self, 'ê²°ê³¼ íŒŒì¼ ì €ìž¥ ê²½ë¡œ ì„ íƒ', self.save_directory)
        
        if directory:
            self.save_directory = directory
            self.save_path_label.setText(f"ì €ìž¥ ê²½ë¡œ: {os.path.basename(directory)}/... ")
        else:
            pass

    # --- ì‹¤ì œ ë³‘í•© ì‹¤í–‰ í•¨ìˆ˜ ---
    def execute_merge(self):
        file_paths = [self.list_widget.item(i).text() for i in range(self.list_widget.count())]
        output_filename = self.output_name_edit.text().strip()
        
        if not file_paths:
            QMessageBox.warning(self, "ê²½ê³ ", "ë³‘í•©í•  PPT íŒŒì¼ì„ 1ê°œ ì´ìƒ ì¶”ê°€í•´ì•¼ í•©ë‹ˆë‹¤.")
            return

        if not output_filename:
            QMessageBox.warning(self, "ê²½ê³ ", "ê²°ê³¼ íŒŒì¼ ì´ë¦„ì„ ìž…ë ¥í•´ì•¼ í•©ë‹ˆë‹¤.")
            return
            
        if not output_filename.lower().endswith(".pptx"):
            output_filename += ".pptx"

        if not self.save_directory:
            QMessageBox.warning(self, "ê²½ê³ ", "ì €ìž¥ ê²½ë¡œë¥¼ ë¨¼ì € ì„ íƒí•´ ì£¼ì„¸ìš”.")
            return

        output_full_path = os.path.join(self.save_directory, output_filename)
        
        self.status_label.setText("ìƒíƒœ: ë³‘í•© ìž‘ì—… ì‹œìž‘ ì¤‘... (PPT íŒŒì¼ ë³€í™˜ í¬í•¨)")
        self.merge_button.setEnabled(False)
        
        self.worker = MergerWorker(file_paths, output_full_path)
        self.worker.merge_finished.connect(self.on_merge_finished)
        self.worker.progress_update.connect(self.on_progress_update)
        self.worker.start()
        
        self.progress_dialog.show()
        
    # --- ì›Œì»¤ ìŠ¤ë ˆë“œ ì‹œê·¸ë„ ì²˜ë¦¬ ---
    def on_progress_update(self, current, total):
        self.progress_dialog.setMaximum(total)
        
        # ë³€í™˜ ë‹¨ê³„ê°€ ìžˆì„ ê²½ìš° ProgressDialog í‘œì‹œ ë³€ê²½
        if total == 1 and current == 1:
            self.progress_dialog.setLabelText("PPT íŒŒì¼ì„ PPTXë¡œ ë³€í™˜ ì¤‘... (MS PowerPoint ì œì–´)")
            self.progress_dialog.setMaximum(100)
            self.progress_dialog.setValue(50)
        else:
            self.progress_dialog.setValue(current)
            self.progress_dialog.setLabelText(f"ìŠ¬ë¼ì´ë“œ ë³µì‚¬ ì¤‘: {current}/{total}")


    def on_merge_finished(self, success, message):
        self.merge_button.setEnabled(True)
        self.progress_dialog.close()
        
        if success:
            self.status_label.setText("ìƒíƒœ: âœ… ë³‘í•© ì™„ë£Œ!")
            QMessageBox.information(self, "ì„±ê³µ", message)
        else:
            self.status_label.setText("ìƒíƒœ: âŒ ì˜¤ë¥˜ ë°œìƒ!")
            QMessageBox.critical(self, "ì˜¤ë¥˜ ë°œìƒ", message)
            
        self.worker = None

# í”„ë¡œê·¸ëž¨ ì‹¤í–‰ ë¶€ë¶„
if __name__ == '__main__':
    app = QApplication(sys.argv)
    
    app.setStyleSheet("""
        QWidget { font-size: 10pt; }
        QPushButton { padding: 8px; border-radius: 5px; }
        QListWidget { border: 1px solid #ddd; padding: 5px; }
        QLabel#status_label { font-weight: bold; padding: 5px; }
    """)
    
    window = PptxMergerApp()
    window.show()
    sys.exit(app.exec_())
